# -*- coding: utf-8 -*-
"""
Created on Sat May  4 15:10:37 2019

@author: willi
"""

def create_charts(prs, dataframes):
    
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.chart import XL_LEGEND_POSITION
    from pptx.util import Inches
    
    df1 = dataframes[0]
    df2 = dataframes[1]
    df3 = dataframes[2]
    df4 = dataframes[3]
    
    # Creating Chart 1
    df1 = df1[['Month', 'Week', 'Hours booked', 'Hours Expected']]

    series_1 = list(df1['Month'])
    series_2 = list(df1['Week'])
    legend = df1.columns[2:4]
    value1 = tuple(df1[legend[0]])
    value2 = tuple(df1[legend[1]])

    # Open ppt and slide
    slide = prs.slides[2]

    #Create the chart
    chart_data = ChartData()
    chart_data.categories = series_2
    chart_data.add_series(legend[0], value1)
    chart_data.add_series(legend[1], value2)

    # Place the chart
    x, y, cx, cy = Inches(1), Inches(2), Inches(6), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    # Create chart 2
    
    series_1 = list(df2.index)
    legend = df2.columns[1:]
    value1 = tuple(df2[legend[0]])
    value2 = tuple(df2[legend[1]])
    value3 = tuple(df2[legend[2]])

    slide = prs.slides[3]

    chart_data = ChartData()
    chart_data.categories = series_1
    chart_data.add_series(legend[0], value1)
    chart_data.add_series(legend[1], value2)
    chart_data.add_series(legend[2], value3)

    x, y, cx, cy = Inches(6), Inches(2), Inches(6), Inches(5)
    chart1 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    # Create chart 3

    series_1 = list(df3.index)
    legend = df3.columns[1:4]
    value1 = tuple(df3[legend[0]])
    value2 = tuple(df3[legend[1]])
    value3 = tuple(df3[legend[2]])

    slide = prs.slides[4]

    chart_data = ChartData()
    chart_data.categories = series_1
    chart_data.add_series(legend[0], value1)
    chart_data.add_series(legend[1], value2)
    chart_data.add_series(legend[2], value3)

    x, y, cx, cy = Inches(0.5), Inches(2.5), Inches(6), Inches(3.5)
    chart1 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # Create chart 4
    series_1 = list(df4.index)
    legend = df4.columns[0:2]
    value1 = tuple(df4[legend[0]])
    value2 = tuple(df4[legend[1]])

    slide = prs.slides[4]

    chart_data = ChartData()
    chart_data.categories = series_1
    chart_data.add_series(legend[0], value1)
    chart_data.add_series(legend[1], value2)


    x, y, cx, cy = Inches(7), Inches(2.5), Inches(6), Inches(3.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart